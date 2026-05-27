from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x13, 0x18, 0x22)   # dark navy
C_CARD      = RGBColor(0x1E, 0x25, 0x33)   # card surface
C_ACCENT    = RGBColor(0x38, 0x8B, 0xFD)   # blue accent
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY_LT   = RGBColor(0xA8, 0xB4, 0xC8)   # secondary text
C_GREY_DIM  = RGBColor(0x55, 0x65, 0x7A)   # muted text
C_LINE      = RGBColor(0x2C, 0x35, 0x46)   # separator
C_LAB       = RGBColor(0x22, 0x4E, 0x7E)   # lab row bg
C_BREAK     = RGBColor(0x1A, 0x20, 0x2C)   # break/lunch row bg
C_REMOTE_BG = RGBColor(0x1A, 0x28, 0x40)   # remote row bg
C_REMOTE_AC = RGBColor(0x2A, 0x6E, 0xC8)   # remote accent strip

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Agenda data ───────────────────────────────────────────────────────────────
DAY1 = [
    ("9:00 – 9:30",   "30 min",  "Welcome, introductions & lab access verification", "Lecture",   True),
    ("9:30 – 10:00",  "30 min",  "SRX Packet walkthrough",                           "Lecture",   True),
    ("10:00 – 12:00", "2 hr",    "Lab 1 · JCL lab exercise",                         "Lab",       True),
    ("12:00 – 1:00",  "60 min",  "Lunch",                                             "Break",     False),
    ("1:00 – 2:00",   "60 min",  "High Availability — MNHA",                         "Lecture",   False),
    ("2:00 – 2:30",   "30 min",  "Break",                                             "Break",     False),
    ("2:30 – 3:30",   "60 min",  "SRX Tap Mode POC",                                 "Lecture",   False),
    ("3:30 – 4:00",   "30 min",  "Break",                                             "Break",     False),
    ("4:00 – 5:00",   "60 min",  "SE Tools",                                          "Lecture",   False),
    ("5:00",          "—",       "Day 1 recap & end of day",                          "Discussion",False),
]

DAY2 = [
    ("9:00 – 9:15",   "15 min",  "Day 1 review",                                     "Discussion",False),
    ("9:15 – 10:15",  "60 min",  "Multi-tenant Support with TSYS and LSYS",          "Lecture",   False),
    ("10:15 – 10:30", "15 min",  "Break",                                             "Break",     False),
    ("10:30 – 11:30", "60 min",  "Sizing and BOM",                                   "Lecture",   False),
    ("11:30 – 12:00", "30 min",  "Sizing and BOM Exercise",                          "Lab",       False),
    ("12:00 – 1:00",  "60 min",  "Lunch",                                             "Break",     False),
    ("1:00 – 2:00",   "60 min",  "Security Director / Security Director Cloud",       "Lecture",   False),
    ("2:00 – 2:30",   "30 min",  "Break",                                             "Break",     False),
    ("2:30 – 3:30",   "60 min",  "Lab · Security Director Cloud Configuration",       "Lab",       False),
    ("3:30 – 4:00",   "30 min",  "Break",                                             "Break",     False),
    ("4:00 – 5:00",   "60 min",  "SRX Automation: MCP Integration",                  "Lecture",   False),
    ("5:00",          "—",       "Day 2 recap & end of day",                          "Discussion",False),
]

DAY3 = [
    ("9:00 – 9:15",   "15 min",  "Day 2 review",                                     "Discussion",False),
    ("9:15 – 10:15",  "60 min",  "Cloud and Virtualization",                         "Lecture",   False),
    ("10:15 – 10:30", "15 min",  "Break",                                             "Break",     False),
    ("10:30 – 11:30", "60 min",  "Advanced Security Threat Lab",                     "Lecture",   False),
    ("11:30 – 12:30", "60 min",  "Lunch",                                             "Break",     False),
    ("1:30 – 2:30",   "60 min",  "EVPN / Secure Data Center Architecture",           "Lecture",   False),
    ("2:30 – 3:00",   "30 min",  "Break",                                             "Break",     False),
    ("3:00 – 3:30",   "30 min",  "CGNAT on MX and SRX",                              "Lecture",   False),
    ("3:30 – 4:15",   "45 min",  "CSDS Scale Out Architecture",                      "Lecture",   False),
    ("4:15 – 4:45",   "30 min",  "Break",                                             "Break",     False),
    ("4:45 – 5:30",   "45 min",  "Wrap-up & next-step learning paths",               "Discussion",False),
]

DAYS = [
    ("Day 01 · Monday",   "JCL Lab, SE Tools & High Availability",        "June 2, 2026",  DAY1),
    ("Day 02 · Tuesday",  "Sizing, Security Director & Automation",        "June 3, 2026",  DAY2),
    ("Day 03 · Wednesday","Cloud, Virtualization & Advanced Security",     "June 4, 2026",  DAY3),
]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=11, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb

# ── Slide 1: Title ────────────────────────────────────────────────────────────
def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title  # may not exist on blank — that's fine

    # Full bg
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)

    # Accent bar left
    add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill=C_ACCENT)

    # Top label
    add_text(slide, "INTERNAL TRAINING  ·  APJ SYSTEMS ENGINEERING",
             Inches(0.4), Inches(0.38), Inches(12), Inches(0.35),
             size=9, color=C_GREY_LT)

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9), Inches(2.2))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Advanced SRX Firewall"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Arial"

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Bootcamp"
    r2.font.size = Pt(44)
    r2.font.bold = True
    r2.font.color.rgb = C_ACCENT
    r2.font.name = "Arial"

    # Subtitle / details
    add_text(slide, "3-Day Intensive  ·  Singapore  ·  June 2 – 4, 2026",
             Inches(0.4), Inches(3.6), Inches(9), Inches(0.4),
             size=14, color=C_GREY_LT)

    # Divider line
    add_rect(slide, Inches(0.4), Inches(4.1), Inches(8.5), Pt(1), fill=C_LINE)

    # Meta pills row
    meta = [
        ("📅", "June 2 – 4, 2026"),
        ("📍", "HPE Singapore · Depot Close"),
        ("👥", "≤ 10 SEs · Internal"),
        ("🏆", "JNCIS-SEC readiness"),
    ]
    for i, (icon, label) in enumerate(meta):
        x = Inches(0.4) + i * Inches(3.1)
        add_rect(slide, x, Inches(4.35), Inches(2.85), Inches(0.52), fill=C_CARD)
        add_text(slide, f"{icon}  {label}",
                 x + Inches(0.12), Inches(4.38), Inches(2.65), Inches(0.45),
                 size=10, color=C_GREY_LT)

    # Footer
    add_text(slide, "CONFIDENTIAL · INTERNAL USE ONLY",
             Inches(0.4), Inches(7.0), Inches(8), Inches(0.3),
             size=8, color=C_GREY_DIM)

# ── Slide 2: Overview ─────────────────────────────────────────────────────────
def make_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)
    add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill=C_ACCENT)

    add_text(slide, "AGENDA OVERVIEW", Inches(0.4), Inches(0.32), Inches(12), Inches(0.3),
             size=9, color=C_ACCENT)
    add_text(slide, "3-Day Schedule at a Glance",
             Inches(0.4), Inches(0.62), Inches(12), Inches(0.55),
             size=22, bold=True, color=C_WHITE)

    add_rect(slide, Inches(0.4), Inches(1.25), Inches(12.5), Pt(1), fill=C_LINE)

    day_summaries = [
        ("Day 01 · Mon · June 2", "JCL Lab, SE Tools & High Availability", [
            "Welcome & lab access verification",
            "SRX Packet walkthrough",
            "Lab 1 · JCL lab exercise",
            "High Availability — MNHA",
            "SRX Tap Mode POC",
            "SE Tools",
        ]),
        ("Day 02 · Tue · June 3", "Sizing, Security Director & Automation", [
            "Multi-tenant Support (TSYS / LSYS)",
            "Sizing and BOM + Exercise",
            "Security Director Cloud",
            "Lab · Security Director Config",
            "SRX Automation: MCP Integration",
        ]),
        ("Day 03 · Wed · June 4", "Cloud, Virtualization & Advanced Security", [
            "Cloud and Virtualization",
            "Advanced Security Threat Lab",
            "EVPN / Secure Data Center Architecture",
            "CGNAT on MX and SRX",
            "CSDS Scale Out Architecture",
            "Wrap-up & next-step learning paths",
        ]),
    ]

    col_w = Inches(4.05)
    for i, (day, theme, sessions) in enumerate(day_summaries):
        x = Inches(0.4) + i * Inches(4.2)
        y_top = Inches(1.4)
        card_h = Inches(5.8)

        # Card bg
        add_rect(slide, x, y_top, col_w, card_h, fill=C_CARD)

        # Accent top strip
        add_rect(slide, x, y_top, col_w, Inches(0.06), fill=C_ACCENT)

        # Day label
        add_text(slide, day, x + Inches(0.18), y_top + Inches(0.14),
                 col_w - Inches(0.3), Inches(0.28),
                 size=9, color=C_ACCENT, bold=True)

        # Theme
        add_text(slide, theme, x + Inches(0.18), y_top + Inches(0.42),
                 col_w - Inches(0.3), Inches(0.55),
                 size=11, bold=True, color=C_WHITE)

        # Divider
        add_rect(slide, x + Inches(0.18), y_top + Inches(1.0),
                 col_w - Inches(0.36), Pt(0.75), fill=C_LINE)

        # Session bullets
        for j, sess in enumerate(sessions):
            sy = y_top + Inches(1.1) + j * Inches(0.68)
            add_rect(slide, x + Inches(0.18), sy + Inches(0.08),
                     Inches(0.06), Inches(0.06), fill=C_ACCENT)
            add_text(slide, sess,
                     x + Inches(0.34), sy, col_w - Inches(0.52), Inches(0.62),
                     size=10, color=C_GREY_LT)

    # Note about remote morning
    add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.28),
             fill=C_REMOTE_BG)
    add_text(slide, "📹  Day 1 morning sessions are remote (via Teams) due to scheduled fire drill at HPE Depot Close.",
             Inches(0.55), Inches(7.12), Inches(12), Inches(0.25),
             size=9, color=C_GREY_LT)

# ── Slide N+2: Day agenda ─────────────────────────────────────────────────────
def make_day_slide(prs, day_label, theme, date, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)
    add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill=C_ACCENT)

    # Header
    add_text(slide, day_label.upper(),
             Inches(0.4), Inches(0.28), Inches(10), Inches(0.28),
             size=9, color=C_ACCENT, bold=True)
    add_text(slide, theme,
             Inches(0.4), Inches(0.55), Inches(10), Inches(0.45),
             size=18, bold=True, color=C_WHITE)
    add_text(slide, date,
             Inches(10.5), Inches(0.28), Inches(2.6), Inches(0.28),
             size=9, color=C_GREY_DIM, align=PP_ALIGN.RIGHT)

    add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Pt(1), fill=C_LINE)

    # Column header
    HDR_Y = Inches(1.12)
    HDR_H = Inches(0.28)
    add_rect(slide, Inches(0.4), HDR_Y, Inches(12.5), HDR_H, fill=C_LINE)
    for label, lx, lw in [
        ("TIME (SGT)", Inches(0.5),  Inches(1.55)),
        ("DUR",        Inches(2.1),  Inches(0.7)),
        ("SESSION",    Inches(2.85), Inches(7.5)),
        ("FORMAT",     Inches(10.5), Inches(2.2)),
    ]:
        add_text(slide, label, lx, HDR_Y + Inches(0.03), lw, HDR_H,
                 size=7.5, color=C_GREY_DIM, bold=True)

    # Rows
    ROW_START_Y = Inches(1.42)
    n = len(rows)
    available_h = Inches(5.85)
    row_h = available_h / n

    for i, (time, dur, session, fmt, remote) in enumerate(rows):
        ry = ROW_START_Y + i * row_h
        is_break = fmt == "Break"
        is_lab   = fmt == "Lab"

        # Row bg
        if remote:
            bg = C_REMOTE_BG
        elif is_break:
            bg = C_BREAK
        elif is_lab:
            bg = C_LAB
        else:
            bg = C_CARD if i % 2 == 0 else C_BG

        add_rect(slide, Inches(0.4), ry, Inches(12.5), row_h - Pt(1), fill=bg)

        # Remote accent strip
        if remote:
            add_rect(slide, Inches(0.4), ry, Inches(0.06), row_h - Pt(1), fill=C_REMOTE_AC)

        txt_color   = C_GREY_DIM if is_break else C_WHITE
        label_color = C_GREY_DIM if is_break else C_GREY_LT
        sess_size   = 9.5 if row_h < Inches(0.55) else 10.5

        add_text(slide, time,    Inches(0.52), ry + Inches(0.04), Inches(1.55), row_h,
                 size=9, color=txt_color, bold=not is_break)
        add_text(slide, dur,     Inches(2.12), ry + Inches(0.04), Inches(0.7),  row_h,
                 size=8.5, color=C_GREY_DIM, italic=True)
        add_text(slide, session, Inches(2.87), ry + Inches(0.04), Inches(7.5),  row_h,
                 size=sess_size, color=txt_color, bold=not is_break)

        # Format badge
        if fmt and not is_break:
            badge_color = C_ACCENT if is_lab else C_GREY_DIM
            badge_bg    = RGBColor(0x1A, 0x35, 0x5A) if is_lab else C_LINE
            add_rect(slide, Inches(10.52), ry + Inches(0.06),
                     Inches(2.0), row_h - Inches(0.14), fill=badge_bg)
            add_text(slide, fmt.upper(),
                     Inches(10.6), ry + Inches(0.07), Inches(1.85), row_h - Inches(0.14),
                     size=7.5, color=badge_color, bold=True, align=PP_ALIGN.CENTER)

        # Remote label
        if remote:
            add_text(slide, "📹 Remote · via Teams",
                     Inches(2.87), ry + row_h - Inches(0.24), Inches(4), Inches(0.22),
                     size=7.5, color=C_REMOTE_AC)

    # Footer
    add_rect(slide, Inches(0.4), Inches(7.22), Inches(12.5), Pt(0.75), fill=C_LINE)
    add_text(slide, "Advanced SRX Firewall Bootcamp · Internal · APJ Systems Engineering",
             Inches(0.4), Inches(7.27), Inches(10), Inches(0.22),
             size=7.5, color=C_GREY_DIM)
    add_text(slide, date, Inches(10.5), Inches(7.27), Inches(2.6), Inches(0.22),
             size=7.5, color=C_GREY_DIM, align=PP_ALIGN.RIGHT)

# ── Build ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

make_title_slide(prs)
make_overview_slide(prs)
for day_label, theme, date, rows in DAYS:
    make_day_slide(prs, day_label, theme, date, rows)

out = "/Users/aleung/Projects/web/srx-bootcamp/SRX_Bootcamp_Agenda.pptx"
prs.save(out)
print(f"Saved: {out}")
